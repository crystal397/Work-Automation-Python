"""
수신자료 파일 점검 도구
- 폴더 내 모든 파일 순회 → 열림 여부 + 손상 여부 점검 → 엑셀 보고서 생성
- 압축 파일(.zip/.7z/.rar) 내부 파일도 개별 점검
- PDF 스캔본/텍스트 구분
- 폴더 깊이(depth)는 실제 구조에서 자동 계산
"""

import os
import sys
import zipfile
import shutil
import tempfile
from pathlib import Path
from datetime import datetime

import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from tqdm import tqdm

# ─────────────────────────────────────────────
# 파일 형식 매핑
# ─────────────────────────────────────────────
EXT_LABEL = {
    ".xlsx": "Microsoft Excel 워크시트",
    ".xls":  "Microsoft Excel 워크시트",
    ".xlsm": "Microsoft Excel 워크시트",
    ".xlsb": "Microsoft Excel 워크시트",
    ".csv":  "CSV 파일",
    ".docx": "Microsoft Word 문서",
    ".doc":  "Microsoft Word 문서",
    ".docm": "Microsoft Word 문서",
    ".pptx": "Microsoft PowerPoint 프레젠테이션",
    ".ppt":  "Microsoft PowerPoint 프레젠테이션",
    ".pptm": "Microsoft PowerPoint 프레젠테이션",
    ".pdf":  "Adobe Acrobat 문서",
    ".hwp":  "한글 문서",
    ".hwpx": "한글 문서",
    ".zip":  "압축(ZIP) 파일",
    ".7z":   "7-Zip 압축 파일",
    ".rar":  "RAR 압축 파일",
    ".tar":  "TAR 압축 파일",
    ".gz":   "GZ 압축 파일",
    ".msg":  "Outlook 항목",
    ".eml":  "이메일 파일",
    ".txt":  "텍스트 파일",
    ".rtf":  "서식 있는 텍스트 파일",
    ".png":  "PNG 이미지",
    ".jpg":  "JPEG 이미지",
    ".jpeg": "JPEG 이미지",
    ".gif":  "GIF 이미지",
    ".bmp":  "BMP 이미지",
    ".tif":  "TIFF 이미지",
    ".tiff": "TIFF 이미지",
    ".dwg":  "AutoCAD 도면",
    ".dxf":  "AutoCAD DXF 도면",
    ".dwf":  "AutoCAD DWF 파일",
    ".mp4":  "MP4 동영상",
    ".avi":  "AVI 동영상",
    ".mov":  "MOV 동영상",
    ".mp3":  "MP3 오디오",
    ".xml":  "XML 파일",
    ".json": "JSON 파일",
    ".mpp":  "Microsoft Project 파일",
    ".log":  "LOG 파일",
    ".bak":  "BAK 파일",
    ".tmp":  "TMP 파일",
    ".dwl":  "DWL 파일",
    ".dwl2": "DWL2 파일",
    ".xer":  "XER 파일",
}


# ─────────────────────────────────────────────
# 파일 열기 시도
# Returns: (비고, 파일형식_override)
#   비고: 정상이면 "", 문제 있으면 사유 문자열
#   파일형식_override: PDF 스캔본/텍스트 등 특수 표기 필요 시, 나머지는 ""
# ─────────────────────────────────────────────
def try_open_file(filepath: Path) -> tuple:
    ext = filepath.suffix.lower()

    try:
        size = filepath.stat().st_size
    except (FileNotFoundError, OSError):
        return ("열기 실패 (경로 오류: 특수문자 등)", "")

    if size == 0:
        return ("파일 손상 (빈 파일)", "")

    try:
        # ── Excel 계열
        if ext in (".xlsx", ".xlsm", ".xls", ".xlsb"):
            # 1차: openpyxl
            try:
                import openpyxl as ox
                wb = ox.load_workbook(_lib_path(filepath), read_only=True, data_only=True)
                has = any(s.max_row and s.max_row > 0 for s in wb.worksheets)
                wb.close()
                return ("", "") if has else ("파일 손상 (내용 없음)", "")
            except Exception:
                pass

            # 2차: xlrd (구형 xls)
            try:
                import xlrd
                wb = xlrd.open_workbook(_lib_path(filepath))
                has = any(wb.sheet_by_index(i).nrows > 0 for i in range(wb.nsheets))
                if has:
                    raw = filepath.read_bytes()
                    irm_kw = [b'\x06DataSpaces', b'EncryptionInfo', b'EncryptedPackage',
                              b'Rights Management', b'\x00I\x00R\x00M',
                              b'aadrm.com', b'Encrypted-Rights-Data',
                              b'rms.microsoft.com', b'_wmcs', b'AUTHENTICATEDDATA']
                    if any(kw in raw for kw in irm_kw):
                        return ("열기 실패 (보안 문서: IRM/DRM 보호)", "")
                return ("", "") if has else ("파일 손상 (내용 없음)", "")
            except Exception:
                pass

            # 3차: 바이너리 분석
            try:
                raw = filepath.read_bytes()
                if raw[:8] == b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1':
                    irm_kw = [b'\x06DataSpaces', b'EncryptionInfo', b'EncryptedPackage',
                              b'Rights Management', b'\x00I\x00R\x00M',
                              b'aadrm.com', b'Encrypted-Rights-Data',
                              b'rms.microsoft.com', b'_wmcs', b'AUTHENTICATEDDATA']
                    if any(kw in raw for kw in irm_kw):
                        return ("열기 실패 (보안 문서: IRM/DRM 보호)", "")
                    return ("", "")
                if b"<html" in raw[:1024].lower() or b"<table" in raw[:1024].lower():
                    return ("", "")
                if len(raw.decode("utf-8", errors="ignore").strip()) > 0:
                    return ("", "")
            except Exception:
                pass

            return ("열기 실패 (형식 불일치 또는 손상)", "")

        # ── Word
        elif ext == ".docx":
            import docx
            doc = docx.Document(_lib_path(filepath))
            all_text = "\n".join(p.text for p in doc.paragraphs)
            if not all_text.strip() and not doc.tables:
                return ("파일 손상 (내용 없음)", "")
            warn = _check_garbled(all_text)
            return (warn, "") if warn else ("", "")

        # ── PowerPoint
        elif ext == ".pptx":
            raw_sig = filepath.read_bytes()[:8]
            # OLE2 시그니처 → 확장자는 .pptx지만 실제 포맷은 구형 .ppt (OLE 기반) 또는 IRM/DRM
            if raw_sig == b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1':
                try:
                    import olefile
                    ole = olefile.OleFileIO(_lib_path(filepath))
                    all_streams = ole.listdir()
                    # IRM/DRM 스트림 감지 (DRMEncryptedDataSpace, EncryptedPackage 등)
                    irm_streams = ("drm", "encrypt", "dataspace", "rights", "_wmcs")
                    for entry in all_streams:
                        entry_str = "/".join(entry).lower()
                        if any(kw in entry_str for kw in irm_streams):
                            ole.close()
                            return ("열기 실패 (보안 문서: IRM/DRM 보호)", "")
                    # 일반 .ppt 포맷 확인
                    has_ppt = ole.exists("PowerPoint Document")
                    ole.close()
                    return ("", "") if has_ppt else ("파일 손상 (PowerPoint Document 스트림 없음)", "")
                except Exception as e:
                    return (f"열기 실패 ({type(e).__name__})", "")
            # 정상 ZIP 기반 .pptx
            from pptx import Presentation
            prs = Presentation(_lib_path(filepath))
            if len(prs.slides) == 0:
                return ("파일 손상 (내용 없음)", "")
            # 슬라이드 전체 텍스트 추출 후 깨짐 검사
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            all_text = "\n".join(texts)
            warn = _check_garbled(all_text)
            return (warn, "") if warn else ("", "")

        # ── PDF
        elif ext == ".pdf":
            import fitz
            doc = fitz.open(_lib_path(filepath))
            if doc.page_count == 0:
                doc.close()
                return ("파일 손상 (내용 없음)", "")
            first_text = doc[0].get_text().lower()
            aip_kw = ["azure information protection", "this is a protected document",
                      "microsoft information protection", "rights management", "irm protected"]
            # 한국형 PDF DRM (Markany, SyncEZ 등) — 1페이지 DRM 안내 메시지 패턴
            krdrm_kw = ["보호된 pdf", "보호된pdf", "뷰어 프로그램을 설치", "뷰어프로그램을 설치",
                        "전용 뷰어", "drm 보호", "drm보호", "문서보안", "문서 보안 솔루션",
                        "syncez", "markany", "fasoo", "irm viewer"]
            if any(kw in first_text for kw in aip_kw + krdrm_kw):
                doc.close()
                return ("열기 실패 (보안 문서: AIP/DRM 보호)", "")
            sample = min(doc.page_count, 5)
            page_texts = [doc[i].get_text() for i in range(sample)]
            total_chars = sum(len(t.strip()) for t in page_texts)
            all_text = "\n".join(page_texts)
            doc.close()
            if total_chars / sample >= 50:
                ftype = "Adobe Acrobat 문서 (텍스트)"
                warn = _check_garbled(all_text)
                return (warn, ftype) if warn else ("", ftype)
            return ("", "Adobe Acrobat 문서 (스캔본)")

        # ── 이메일
        elif ext in (".msg", ".eml"):
            return ("", "") if size > 100 else ("파일 손상 (내용 없음)", "")

        # ── Word (구형 .doc) — olefile
        elif ext == ".doc":
            try:
                import olefile
                if not olefile.isOleFile(_lib_path(filepath)):
                    return ("파일 손상 (OLE 구조 오류)", "")
                ole = olefile.OleFileIO(_lib_path(filepath))
                all_streams = ole.listdir()
                # IRM/DRM 스트림 감지
                irm_streams = ("drm", "encrypt", "dataspace", "rights", "_wmcs")
                for entry in all_streams:
                    if any(kw in "/".join(entry).lower() for kw in irm_streams):
                        ole.close()
                        return ("열기 실패 (보안 문서: IRM/DRM 보호)", "")
                if not ole.exists("WordDocument"):
                    ole.close()
                    return ("파일 손상 (WordDocument 스트림 없음)", "")
                ole.close()
                return ("", "")
            except Exception as e:
                return (f"열기 실패 ({type(e).__name__})", "")

        # ── PowerPoint (구형 .ppt) — olefile
        elif ext == ".ppt":
            try:
                import olefile
                if not olefile.isOleFile(_lib_path(filepath)):
                    return ("파일 손상 (OLE 구조 오류)", "")
                ole = olefile.OleFileIO(_lib_path(filepath))
                all_streams = ole.listdir()
                # IRM/DRM 스트림 감지
                irm_streams = ("drm", "encrypt", "dataspace", "rights", "_wmcs")
                for entry in all_streams:
                    if any(kw in "/".join(entry).lower() for kw in irm_streams):
                        ole.close()
                        return ("열기 실패 (보안 문서: IRM/DRM 보호)", "")
                if not ole.exists("PowerPoint Document"):
                    ole.close()
                    return ("파일 손상 (PowerPoint Document 스트림 없음)", "")
                ole.close()
                return ("", "")
            except Exception as e:
                return (f"열기 실패 ({type(e).__name__})", "")

        # ── 한글 문서
        elif ext in (".hwp", ".hwpx"):
            try:
                raw = filepath.read_bytes()
                # Microsoft RMS/AIP .pfile 컨테이너 (HWP + AIP 암호화)
                if raw[:6] == b'.pfile':
                    return ("열기 실패 (보안 문서: AIP/DRM 보호)", "")
                # HWP 5.x: OLE2 기반 (D0CF11E0 시그니처)
                if raw[:8] == b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1':
                    import olefile
                    if not olefile.isOleFile(_lib_path(filepath)):
                        return ("파일 손상 (OLE 구조 오류)", "")
                    ole = olefile.OleFileIO(_lib_path(filepath))
                    # 암호화 스트림 확인
                    all_streams = ole.listdir()
                    enc_keywords = ("encrypt", "security", "drm", "protect")
                    for entry in all_streams:
                        entry_str = "/".join(entry).lower()
                        if any(kw in entry_str for kw in enc_keywords):
                            ole.close()
                            return ("열기 실패 (암호 보호)", "")
                    # HWP 핵심 스트림 존재 확인
                    if not ole.exists("FileHeader"):
                        ole.close()
                        return ("파일 손상 (HWP FileHeader 없음)", "")
                    # BodyText 본문 스트림 검사
                    if not ole.exists("BodyText/Section0"):
                        ole.close()
                        return ("파일 손상 (HWP 본문 없음)", "")
                    raw_section = ole.openstream("BodyText/Section0").read()
                    ole.close()
                    # ParaText 레코드 파싱으로 실제 텍스트 추출 후 깨짐 검사
                    hwp_text = _extract_hwp_text(raw_section)
                    warn = _check_garbled(hwp_text)
                    if warn:
                        return (warn, "")
                # HWPX: ZIP 기반 (PK 시그니처)
                elif raw[:2] == b'PK':
                    import xml.etree.ElementTree as _ET
                    if not zipfile.is_zipfile(str(filepath)):
                        return ("파일 손상 (ZIP 구조 오류)", "")
                    with zipfile.ZipFile(str(filepath)) as z:
                        section_files = sorted(
                            n for n in z.namelist()
                            if n.startswith("Contents/section") and n.endswith(".xml")
                        )
                        if not section_files:
                            return ("파일 손상 (HWPX 본문 없음)", "")
                        texts = []
                        for sname in section_files[:3]:
                            try:
                                xml_bytes = z.read(sname)
                                root = _ET.fromstring(xml_bytes)
                                for elem in root.iter():
                                    if elem.tag.endswith("}t") or elem.tag == "t":
                                        if elem.text:
                                            texts.append(elem.text)
                            except _ET.ParseError:
                                return ("파일 손상 (HWPX XML 오류)", "")
                    all_text = "\n".join(texts)
                    warn = _check_garbled(all_text)
                    if warn:
                        return (warn, "")
                return ("", "")
            except Exception as e:
                return (f"열기 실패 ({type(e).__name__})", "")

        # ── 이미지 — Pillow로 유효성 확인 및 해상도 표시
        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff"):
            try:
                from PIL import Image
                with Image.open(filepath) as img:
                    img.verify()
                with Image.open(filepath) as img:
                    w, h = img.size
                    mode = img.mode
                return ("", f"{ext[1:].upper()} 이미지 ({w}×{h}, {mode})")
            except ImportError:
                return ("", "")  # Pillow 미설치
            except Exception:
                return ("파일 손상 (이미지 오류)", "")

        # ── AutoCAD DXF — ezdxf로 파싱, 없으면 텍스트 시그니처 확인
        elif ext == ".dxf":
            try:
                import ezdxf
                doc = ezdxf.readfile(_lib_path(filepath))
                count = len(list(doc.modelspace()))
                return ("", f"AutoCAD DXF 도면 ({count:,}개 객체)")
            except ImportError:
                try:
                    text = filepath.read_text(encoding="utf-8", errors="replace")
                    if "SECTION" in text and "ENDSEC" in text:
                        return ("", "")
                    return ("파일 손상 (DXF 구조 오류)", "")
                except Exception:
                    return ("열기 실패", "")
            except Exception as e:
                return (f"열기 실패 ({type(e).__name__})", "")

        # ── AutoCAD DWG — 파일 시그니처(AC####) 확인
        elif ext == ".dwg":
            try:
                sig = filepath.read_bytes()[:6]
                if sig[:2] == b"AC":
                    ver = sig.decode("ascii", errors="replace")
                    return ("", f"AutoCAD 도면 ({ver})")
                return ("파일 손상 (DWG 시그니처 오류)", "")
            except Exception as e:
                return (f"열기 실패 ({type(e).__name__})", "")

        # ── DWL / DWL2 — AutoCAD 잠금 파일 (텍스트)
        elif ext in (".dwl", ".dwl2"):
            try:
                text = filepath.read_bytes().decode("utf-8", errors="replace")
                return ("", "") if text.strip() else ("파일 손상 (내용 없음)", "")
            except Exception as e:
                return (f"열기 실패 ({type(e).__name__})", "")

        # ── XER — Primavera P6 (ERMHDR 시그니처 또는 XML)
        elif ext == ".xer":
            try:
                head = filepath.read_bytes()[:200].decode("utf-8", errors="replace")
                if "ERMHDR" in head:
                    return ("", "")
                import xml.etree.ElementTree as ET
                ET.parse(str(filepath))
                return ("", "")
            except Exception:
                try:
                    if filepath.read_text(encoding="utf-8", errors="replace").strip():
                        return ("", "")
                except Exception:
                    pass
                return ("파일 손상 (XER 구조 오류)", "")

        # ── 텍스트 계열 — 인코딩 순차 시도
        elif ext in (".txt", ".log", ".csv", ".rtf"):
            for enc in ("utf-8-sig", "utf-8", "cp949", "utf-16"):
                try:
                    text = filepath.read_text(encoding=enc)
                    return ("", "") if text.strip() else ("파일 손상 (내용 없음)", "")
                except (UnicodeDecodeError, UnicodeError):
                    continue
            return ("열기 실패 (인코딩 불명)", "")

        # ── 동영상 / 오디오 — 파일 시그니처 확인
        elif ext in (".mp4", ".mov", ".avi", ".mp3"):
            try:
                raw = filepath.read_bytes()[:12]
                if ext in (".mp4", ".mov") and raw[4:8] in (b"ftyp", b"moov", b"mdat", b"wide", b"free"):
                    return ("", "")
                if ext == ".avi" and raw[:4] == b"RIFF" and raw[8:12] == b"AVI ":
                    return ("", "")
                if ext == ".mp3" and (raw[:3] == b"ID3" or raw[:2] == b"\xff\xfb"):
                    return ("", "")
                # 시그니처 불일치 — 크기는 있으므로 경고만
                return ("열기 실패 (시그니처 불일치)", "")
            except Exception as e:
                return (f"열기 실패 ({type(e).__name__})", "")

        # ── BAK / TMP — 시그니처로 원본 형식 추정
        elif ext in (".bak", ".tmp"):
            try:
                raw = filepath.read_bytes()[:8]
                if raw[:2] == b"PK":
                    return ("", "ZIP 기반 백업 (docx/xlsx 등)")
                if raw[:8] == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":
                    return ("", "OLE 기반 백업 (doc/xls 등)")
                if raw[:4] == b"%PDF":
                    return ("", "PDF 기반 백업")
                for enc in ("utf-8-sig", "utf-8", "cp949"):
                    try:
                        text = filepath.read_text(encoding=enc)
                        if text.strip():
                            return ("", "텍스트 기반 백업")
                    except (UnicodeDecodeError, UnicodeError):
                        continue
                return ("", "바이너리 백업")
            except Exception as e:
                return (f"열기 실패 ({type(e).__name__})", "")

        # ── 그 외: 크기로만 판단
        else:
            return ("", "")

    except PermissionError:
        return ("열기 실패 (접근 권한 없음)", "")
    except zipfile.BadZipFile:
        return ("파일 손상 (압축 파일 오류)", "")
    except Exception as e:
        msg = str(e).lower()
        if "password" in msg or "encrypt" in msg or "암호" in msg:
            return ("열기 실패 (암호 보호)", "")
        elif "permission" in msg or "access" in msg or "denied" in msg:
            return ("열기 실패 (접근 권한 없음)", "")
        elif "login" in msg or "auth" in msg or "credential" in msg:
            return ("열기 실패 (로그인 필요)", "")
        elif "corrupt" in msg or "invalid" in msg or "truncat" in msg:
            return ("파일 손상 (파일 오류)", "")
        elif type(e).__name__ in ("FileDataError",):
            return ("파일 손상 또는 보호된 파일 (열기 실패)", "")
        else:
            return (f"열기 실패 ({type(e).__name__})", "")


# ─────────────────────────────────────────────
# 압축 해제
# ─────────────────────────────────────────────
def extract_archive(filepath: Path, dest: Path) -> tuple:
    """압축 해제. (성공여부, 오류메시지) 반환."""
    ext = filepath.suffix.lower()
    try:
        if ext == ".zip":
            with zipfile.ZipFile(filepath, "r") as z:
                for info in z.infolist():
                    if not (info.flag_bits & 0x800):
                        try:
                            info.filename = info.filename.encode("cp437").decode("cp949")
                        except Exception:
                            try:
                                info.filename = info.filename.encode("cp437").decode("utf-8")
                            except Exception:
                                pass
                    info.filename = info.filename.replace("\\", "/")
                    z.extract(info, dest)
            return (True, "")
        elif ext == ".7z":
            import py7zr
            with py7zr.SevenZipFile(filepath, mode="r") as z:
                z.extractall(dest)
            return (True, "")
        elif ext == ".rar":
            import rarfile
            with rarfile.RarFile(filepath, "r") as z:
                z.extractall(dest)
            return (True, "")
        else:
            return (False, "지원하지 않는 압축 형식")
    except Exception as e:
        msg = str(e).lower()
        if "password" in msg or "encrypt" in msg:
            return (False, "열기 실패 (암호 보호)")
        return (False, f"압축 해제 실패 ({type(e).__name__})")


# ─────────────────────────────────────────────
# 유틸
# ─────────────────────────────────────────────
def to_extended_path(p: Path) -> Path:
    """Windows 긴 경로 처리 (UNC 확장 경로 적용)"""
    s = str(p)
    prefix = "\\\\?\\"
    if os.name == "nt" and not s.startswith(prefix):
        s = prefix + s
    return Path(s)


def _lib_path(p: Path) -> str:
    r"""라이브러리에 넘길 경로 문자열: \\?\ 접두사 제거 (openpyxl·fitz 등 C 확장 호환)"""
    s = str(p)
    if s.startswith("\\\\?\\"):
        s = s[4:]
    return s


def _extract_hwp_text(raw_section: bytes) -> str:
    """
    HWP 5.x BodyText/Section0 스트림에서 텍스트를 추출한다.
    - zlib raw deflate 압축 해제 후 HWP 레코드를 순회
    - ParaText 레코드 (태그 ID = 67) 의 데이터를 UTF-16LE 로 디코딩
    """
    import zlib
    try:
        data = zlib.decompress(raw_section, -15)
    except zlib.error:
        data = raw_section  # 비압축 모드

    TAG_PARA_TEXT = 67
    texts = []
    pos = 0
    while pos + 4 <= len(data):
        header = int.from_bytes(data[pos:pos + 4], "little")
        tag_id = header & 0x3FF
        size   = (header >> 20) & 0xFFF
        pos += 4
        if size == 0xFFF:            # 확장 크기 필드
            if pos + 4 > len(data):
                break
            size = int.from_bytes(data[pos:pos + 4], "little")
            pos += 4
        if pos + size > len(data):
            break
        if tag_id == TAG_PARA_TEXT and size > 0:
            try:
                texts.append(data[pos:pos + size].decode("utf-16-le", errors="replace"))
            except Exception:
                pass
        pos += size
    return "\n".join(texts)


def _garbling_ratio(text: str) -> float:
    """
    텍스트에서 깨진 문자 비율을 반환 (0.0 ~ 1.0).
    - U+FFFD (대체문자): 인코딩 변환 실패 시 생성
    - 제어문자 (탭·개행·CR 제외): 바이너리가 텍스트 필드에 흘러든 경우
    """
    if not text:
        return 0.0
    bad = sum(
        1 for c in text
        if c == "\ufffd" or (ord(c) < 32 and c not in "\t\n\r")
    )
    return bad / len(text)


def _check_garbled(text: str, min_len: int = 30, threshold: float = 0.15) -> str:
    """
    텍스트 깨짐 여부를 검사.
    - min_len 미만이면 판단 보류 (빈 문서는 별도 처리)
    - 깨짐 비율이 threshold 이상이면 경고 문자열 반환, 정상이면 ""
    """
    text = text.strip()
    if len(text) < min_len:
        return ""
    ratio = _garbling_ratio(text)
    if ratio >= threshold:
        return f"내용 깨짐 의심 (깨진 문자 {ratio:.0%})"
    return ""


def ext_to_label(ext: str) -> str:
    return EXT_LABEL.get(ext, ext[1:].upper() + " 파일" if ext else "알 수 없음")


def get_folder_hierarchy(root: Path, filepath: Path, max_depth: int) -> list:
    """루트 기준 상대 경로에서 폴더 계층 추출 (max_depth 단계)"""
    rel = filepath.relative_to(root)
    parts = list(rel.parts[:-1])
    while len(parts) < max_depth:
        parts.append("")
    return parts[:max_depth]


def _build_archive_parts(raw_fs_parts: list, zip_name: str,
                         inner_folder_parts: list, max_depth: int) -> list:
    """
    압축 내부 항목의 계층 구성.
    raw_fs_parts   : zip 파일이 있는 폴더까지의 경로 (패딩 없음)
    zip_name       : zip 파일명 → 하나의 depth 열
    inner_folder_parts : zip 내부 폴더 경로 부분 (패딩 없음)
    """
    parts = raw_fs_parts + [zip_name] + inner_folder_parts
    while len(parts) < max_depth:
        parts.append("")
    return parts[:max_depth]


def make_rec(parts: list, fname: str, ftype: str, note: str, path_str: str) -> dict:
    rec = {f"{i+1}단계": parts[i] for i in range(len(parts))}
    rec.update({
        "파일명": fname,
        "파일형식": ftype,
        "비고": note,
        "_path": path_str,
        "_flagged": bool(note),
    })
    return rec


# ─────────────────────────────────────────────
# 압축 내부 점검
# ─────────────────────────────────────────────
ARCHIVE_EXTS = {".zip", ".7z", ".rar"}


def scan_archive(filepath: Path, raw_fs_parts: list, zip_name: str, max_depth: int) -> list:
    """
    raw_fs_parts: zip 파일이 있는 폴더까지의 경로 (패딩 없음)
    zip 내부 폴더 구조를 depth 열에 직접 배분.
    """
    records = []
    tmp_dir = Path(tempfile.mkdtemp(prefix="chk_"))
    try:
        ok, err = extract_archive(filepath, tmp_dir)
        if not ok:
            parts = _build_archive_parts(raw_fs_parts, zip_name, [], max_depth)
            records.append(make_rec(parts, zip_name, ext_to_label(filepath.suffix.lower()), err, str(filepath)))
            return records

        inner_files = sorted(f for f in tmp_dir.rglob("*") if f.is_file())
        if not inner_files:
            parts = _build_archive_parts(raw_fs_parts, zip_name, [], max_depth)
            records.append(make_rec(parts, "", "", "파일 손상 (빈 압축파일)", str(filepath)))
            return records

        for inner in inner_files:
            inner_rel   = inner.relative_to(tmp_dir)
            inner_parts = list(inner_rel.parts[:-1])   # 내부 폴더 경로
            fname       = inner_rel.name               # 파일명만

            inner_ext = inner.suffix.lower()
            note, ftype_override = try_open_file(inner)
            ftype = ftype_override if ftype_override else ext_to_label(inner_ext)
            parts = _build_archive_parts(raw_fs_parts, zip_name, inner_parts, max_depth)
            records.append(make_rec(parts, fname, ftype, note,
                                    str(filepath) + "/" + str(inner_rel)))
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)
    return records


# ─────────────────────────────────────────────
# 폴더 순회
# ─────────────────────────────────────────────
def _peek_archive_depth(filepath: Path) -> int:
    """압축 파일 내부의 최대 폴더 깊이 반환 (실패 시 0)"""
    ext = filepath.suffix.lower()
    try:
        if ext == ".zip":
            with zipfile.ZipFile(filepath, "r") as z:
                depths = [
                    len(Path(name.replace("\\", "/")).parts) - 1
                    for name in z.namelist()
                    if not name.endswith("/")
                ]
                return max(depths, default=0)
        elif ext == ".7z":
            import py7zr
            with py7zr.SevenZipFile(filepath, mode="r") as z:
                depths = [len(Path(n).parts) - 1 for n in z.getnames()
                          if not n.endswith("/")]
                return max(depths, default=0)
        elif ext == ".rar":
            import rarfile
            with rarfile.RarFile(filepath, "r") as rf:
                depths = [len(Path(i.filename).parts) - 1
                          for i in rf.infolist() if not i.is_dir()]
                return max(depths, default=0)
    except Exception:
        pass
    return 0


def _compute_max_depth(root: Path) -> int:
    """
    실제 폴더 + 압축 내부 구조까지 반영한 최대 깊이 계산 (최소 1).
    일반 파일: 파일 위치까지의 폴더 depth
    압축 파일: (파일 위치 depth) + 1(zip명) + (내부 폴더 depth)
    """
    # 먼저 모든 파일 목록 수집
    all_files = []
    for dirpath, _, filenames in os.walk(root):
        fs_depth = len(Path(dirpath).relative_to(root).parts)
        for fname in filenames:
            all_files.append((Path(dirpath) / fname, fs_depth))

    max_d = 1
    archive_files = [(fp, d) for fp, d in all_files if fp.suffix.lower() in ARCHIVE_EXTS]
    normal_files  = [(fp, d) for fp, d in all_files if fp.suffix.lower() not in ARCHIVE_EXTS]

    # 일반 파일
    for _, fs_depth in normal_files:
        max_d = max(max_d, fs_depth)

    # 압축 파일 — 내부 depth까지 합산 (진행바 표시)
    if archive_files:
        for fp, fs_depth in tqdm(archive_files, desc="압축 내부 depth 분석", unit="개"):
            internal = _peek_archive_depth(fp)
            max_d = max(max_d, fs_depth + 1 + internal)

    return max(max_d, 1)


def scan_folder(root: Path) -> tuple[list, int]:
    """
    폴더 전체를 순회하여 파일 점검 결과를 반환.
    반환: (records, max_depth)
    """
    records = []
    root = to_extended_path(root.resolve())

    # 실제 폴더 깊이 계산
    max_depth = _compute_max_depth(root)
    print(f"[INFO] 폴더 최대 깊이: {max_depth}단계")

    # 전체 파일 목록 수집 (tqdm 전체 카운트 확보)
    all_entries = []  # (dir_path, fname | None)  None이면 빈 폴더
    for dirpath, dirnames, filenames in os.walk(root):
        dirnames.sort()
        dir_path = Path(dirpath)
        if not filenames:
            has_child = any(files for _, _, files in os.walk(dirpath))
            if not has_child:
                all_entries.append((dir_path, None))
        else:
            for fname in sorted(filenames):
                all_entries.append((dir_path, fname))

    # 점검
    with tqdm(all_entries, desc="파일 점검", unit="개", dynamic_ncols=True) as bar:
        for dir_path, fname in bar:
            if fname is None:
                parts = get_folder_hierarchy(root, dir_path / "_dummy_", max_depth)
                rec = make_rec(parts, "", "", "폴더 비었음", str(dir_path))
                rec["_flagged"] = False
                records.append(rec)
                continue

            bar.set_postfix(파일=fname[:30])
            fpath = dir_path / fname
            ext = fpath.suffix.lower()
            parts = get_folder_hierarchy(root, fpath, max_depth)

            if ext in ARCHIVE_EXTS:
                raw_fs_parts = list(fpath.relative_to(root).parts[:-1])
                records.extend(scan_archive(fpath, raw_fs_parts, fname, max_depth))
            else:
                note, ftype_override = try_open_file(fpath)
                ftype = ftype_override if ftype_override else ext_to_label(ext)
                records.append(make_rec(parts, fname, ftype, note, str(fpath)))

    return records, max_depth


# ─────────────────────────────────────────────
# 스타일 상수
# ─────────────────────────────────────────────
HEADER_FILL = PatternFill("solid", fgColor="1F4E79")
DAMAGE_FILL = PatternFill("solid", fgColor="FFE0E0")
EMPTY_FILL  = PatternFill("solid", fgColor="FFF2CC")
ALT_FILL    = PatternFill("solid", fgColor="F2F7FF")
WHITE_FILL  = PatternFill("solid", fgColor="FFFFFF")
THIN        = Side(style="thin", color="AAAAAA")
THIN_BORDER = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)
HEADER_FONT = Font(name="맑은 고딕", size=10, bold=True, color="FFFFFF")
NORMAL_FONT = Font(name="맑은 고딕", size=9)
BOLD_FONT   = Font(name="맑은 고딕", size=9, bold=True)
DAMAGE_FONT = Font(name="맑은 고딕", size=9, color="C00000", bold=True)
TITLE_FONT  = Font(name="맑은 고딕", size=14, bold=True, color="FFFFFF")
CENTER = Alignment(horizontal="center", vertical="center", wrap_text=True)
LEFT   = Alignment(horizontal="left",   vertical="center", wrap_text=True)


def write_cell(ws, row, col, value, font=None, fill=None, alignment=None, border=None):
    c = ws.cell(row=row, column=col, value=value)
    if font:      c.font      = font
    if fill:      c.fill      = fill
    if alignment: c.alignment = alignment
    if border:    c.border    = border
    return c


# ─────────────────────────────────────────────
# 엑셀 보고서 생성
# ─────────────────────────────────────────────
def build_excel(records: list, out_path: str, folder_path: str, max_depth: int):
    # depth에 따라 동적으로 계산
    total_cols    = 1 + max_depth + 3          # 번호 + N단계 + 파일명 + 파일형식 + 비고
    last_col_letter = get_column_letter(total_cols)
    headers       = ["번호"] + [f"{i+1}단계" for i in range(max_depth)] + ["파일명", "파일형식", "비고"]

    # 열 너비: 번호(5) | 단계열(20 each) | 파일명(40) | 파일형식(20) | 비고(28)
    col_widths = [5] + [20] * max_depth + [40, 20, 28]

    # 파일명·파일형식·비고 열 인덱스 (1-based)
    fname_col  = 1 + max_depth + 1
    ftype_col  = 1 + max_depth + 2
    note_col   = 1 + max_depth + 3

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "수신자료 정리"

    for col_idx, width in enumerate(col_widths, 1):
        ws.column_dimensions[get_column_letter(col_idx)].width = width

    # 타이틀
    ws.row_dimensions[1].height = 28
    ws.merge_cells(f"A1:{last_col_letter}1")
    c = ws["A1"]
    c.value     = f"수신자료 정리 / {datetime.now().strftime('%y%m%d')}_파일 수령 현황"
    c.font      = TITLE_FONT
    c.fill      = HEADER_FILL
    c.alignment = CENTER

    # 경로
    ws.row_dimensions[2].height = 18
    ws.merge_cells(f"A2:{last_col_letter}2")
    c = ws["A2"]
    c.value     = f"대상 폴더: {folder_path}"
    c.font      = Font(name="맑은 고딕", size=8, italic=True, color="595959")
    c.fill      = PatternFill("solid", fgColor="D6E4F0")
    c.alignment = LEFT

    # 헤더
    ws.row_dimensions[3].height = 22
    for col_idx, hdr in enumerate(headers, 1):
        write_cell(ws, 3, col_idx, hdr, font=HEADER_FONT, fill=HEADER_FILL,
                   alignment=CENTER, border=THIN_BORDER)

    # 데이터
    flagged_list = []
    center_cols = {1, fname_col, ftype_col, note_col}

    for row_num, (idx, rec) in tqdm(
        enumerate(enumerate(records, 1), 4),
        total=len(records), desc="엑셀 작성", unit="행"
    ):
        ws.row_dimensions[row_num].height = 16
        is_flagged = rec["_flagged"]
        bigo = rec["비고"]
        is_empty = "폴더 비었음" in bigo

        row_fill = (
            EMPTY_FILL  if is_empty  else
            DAMAGE_FILL if is_flagged else
            ALT_FILL    if row_num % 2 == 0 else
            WHITE_FILL
        )

        vals = ([idx]
                + [rec[f"{i+1}단계"] for i in range(max_depth)]
                + [rec["파일명"], rec["파일형식"], bigo])

        for col_idx, val in enumerate(vals, 1):
            font = (DAMAGE_FONT if is_flagged and col_idx in (fname_col, note_col)
                    else BOLD_FONT if col_idx == 1
                    else NORMAL_FONT)
            write_cell(ws, row_num, col_idx, val,
                       font=font, fill=row_fill,
                       alignment=CENTER if col_idx in center_cols else LEFT,
                       border=THIN_BORDER)

        if is_flagged and not is_empty:
            flagged_list.append(rec)

    # 집계
    mid = total_cols // 2
    summary_row = 3 + len(records) + 2
    ws.row_dimensions[summary_row].height = 20

    ws.merge_cells(f"A{summary_row}:{get_column_letter(mid)}{summary_row}")
    c = ws.cell(row=summary_row, column=1,
                value=f"총 파일 수: {len([r for r in records if r['파일명']])}건")
    c.font = BOLD_FONT
    c.fill = PatternFill("solid", fgColor="D6E4F0")
    c.alignment = CENTER
    c.border = THIN_BORDER

    start_col = mid + 1
    ws.merge_cells(f"{get_column_letter(start_col)}{summary_row}:{last_col_letter}{summary_row}")
    c = ws.cell(row=summary_row, column=start_col, value=f"비고 기록: {len(flagged_list)}건")
    c.font = Font(name="맑은 고딕", size=9, bold=True, color="C00000")
    c.fill = DAMAGE_FILL
    c.alignment = CENTER
    c.border = THIN_BORDER

    # 비고 목록 시트
    ws2 = wb.create_sheet("비고 기록 파일 목록")
    for col_idx, width in enumerate(col_widths, 1):
        ws2.column_dimensions[get_column_letter(col_idx)].width = width

    ws2.row_dimensions[1].height = 28
    ws2.merge_cells(f"A1:{last_col_letter}1")
    c = ws2["A1"]
    c.value     = "비고 기록 파일 목록"
    c.font      = TITLE_FONT
    c.fill      = PatternFill("solid", fgColor="C00000")
    c.alignment = CENTER

    ws2.row_dimensions[2].height = 22
    for col_idx, hdr in enumerate(headers, 1):
        write_cell(ws2, 2, col_idx, hdr, font=HEADER_FONT,
                   fill=PatternFill("solid", fgColor="C00000"),
                   alignment=CENTER, border=THIN_BORDER)

    for r_idx, rec in enumerate(flagged_list, 3):
        ws2.row_dimensions[r_idx].height = 16
        vals = ([r_idx - 2]
                + [rec[f"{i+1}단계"] for i in range(max_depth)]
                + [rec["파일명"], rec["파일형식"], rec["비고"]])
        for col_idx, val in enumerate(vals, 1):
            write_cell(ws2, r_idx, col_idx, val,
                       font=DAMAGE_FONT if col_idx in (fname_col, note_col) else NORMAL_FONT,
                       fill=DAMAGE_FILL,
                       alignment=CENTER if col_idx in center_cols else LEFT,
                       border=THIN_BORDER)

    wb.save(out_path)
    return flagged_list


# ─────────────────────────────────────────────
# 경로 입력 (GUI → 콘솔 fallback)
# ─────────────────────────────────────────────
def get_folder_path_gui() -> str:
    import tkinter as tk
    from tkinter import filedialog, messagebox
    root = tk.Tk()
    root.withdraw()
    messagebox.showinfo("수신자료 파일 점검", "점검할 폴더를 선택해 주세요.")
    folder = filedialog.askdirectory(title="점검 대상 폴더 선택")
    root.destroy()
    return folder


def get_folder_path_console() -> str:
    print("=" * 50)
    print("  수신자료 파일 점검 도구")
    print("=" * 50)
    print()
    while True:
        path = input("점검할 폴더 경로를 입력하세요: ").strip().strip('"')
        if path:
            return path
        print("[ERROR] 경로를 입력해 주세요.")


# ─────────────────────────────────────────────
# 실행
# ─────────────────────────────────────────────
if __name__ == "__main__":
    folder_path = ""
    try:
        folder_path = get_folder_path_gui()
        if not folder_path:
            print("[INFO] GUI 취소 → 콘솔 입력으로 전환합니다.")
            folder_path = get_folder_path_console()
    except Exception as e:
        print(f"[INFO] GUI 사용 불가 ({e}) → 콘솔 입력으로 전환합니다.")
        folder_path = get_folder_path_console()

    folder = Path(folder_path)
    if not folder.exists():
        print(f"[ERROR] 폴더를 찾을 수 없습니다: {folder_path}")
        input("엔터를 눌러 종료...")
        sys.exit(1)

    print(f"[INFO] 폴더 순회 시작: {folder_path}")
    records, max_depth = scan_folder(folder)
    print(f"[INFO] 총 {len(records)}건 확인 완료")

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    excel_out = f"수신자료_파일점검결과_{timestamp}.xlsx"

    flagged = build_excel(records, excel_out, folder_path, max_depth)
    print(f"[OK] 엑셀 보고서 저장: {excel_out}")
    print(f"[INFO] 비고 기록 파일: {len(flagged)}건")

    print("\n=== 완료 ===")
    print(f"  보고서: {excel_out}")
    if flagged:
        print(f"\n  비고 기록 파일 목록 ({len(flagged)}건):")
        for rec in flagged:
            cat = " > ".join(
                rec[f"{i+1}단계"] for i in range(max_depth) if rec[f"{i+1}단계"]
            )
            print(f"    [{cat}]  {rec['파일명']}  →  {rec['비고']}")

    input("\n엔터를 눌러 종료...")
